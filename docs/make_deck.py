#!/usr/bin/env python3
"""Build Artha-Pitch.pptx — Verge-themed deck embedding the app screenshots.
Run:  /tmp/pptenv/bin/python docs/make_deck.py"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
SHOTS = os.path.join(HERE, "screenshots")

CANVAS = RGBColor(0x13, 0x13, 0x13)
MINT   = RGBColor(0x3C, 0xFF, 0xD0)
UV     = RGBColor(0x5C, 0x33, 0xFF)   # lightened ultraviolet for contrast on black
YELLOW = RGBColor(0xFF, 0xD8, 0x4D)
PINK   = RGBColor(0xFF, 0x3B, 0x6B)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
MUTED  = RGBColor(0x94, 0x94, 0x94)
LINE   = RGBColor(0x31, 0x31, 0x31)

DISPLAY = "Arial Black"
BODY    = "Arial"
MONO    = "Consolas"

PHONE_RATIO = 1440 / 3168  # width/height

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = CANVAS
    return s


def rule(s, x, y, w, color=UV, h=0.045):
    shp = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background(); shp.shadow.inherit = False
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=4):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, line in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space)
        for (t, sz, col, font, bold) in line:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col
            r.font.name = font; r.font.bold = bold
    return tb


def kicker(s, x, y, t, color=MINT):
    text(s, x, y, 11, 0.4, [[(t.upper(), 13, color, MONO, True)]])


def phone(s, name, x, y, h):
    w = h * PHONE_RATIO
    s.shapes.add_picture(os.path.join(SHOTS, name), Inches(x), Inches(y), height=Inches(h))
    return w


# ── 1. Title ────────────────────────────────────────────────────────────────
s = slide()
rule(s, 0.9, 2.55, 4.2, UV)
text(s, 0.85, 1.4, 11, 1.4, [[("ARTHA", 88, WHITE, DISPLAY, True)]])
text(s, 0.9, 2.8, 11.5, 1.0,
     [[("Turning the kirana counter into a ", 26, WHITE, BODY, False),
       ("financial identity", 26, MINT, BODY, True),
       (" for the next 500 million.", 26, WHITE, BODY, False)]])
text(s, 0.9, 4.2, 11.7, 1.4,
     [[("An AI commerce-and-finance OS for India's 13 million corner shops — ", 17, MUTED, BODY, False),
       ("the phone is the only tool.", 17, MUTED, BODY, True)],
      [("Local-first — runs on the iQOO's on-device edge; the cloud is touched only to OCR a photo.", 17, MINT, BODY, True)]])
text(s, 0.9, 6.55, 12, 0.6,
     [[("PS 3 · REIMAGINE MONEY FOR BHARAT", 13, MINT, MONO, True),
       ("      iQOO Hackathon", 13, MUTED, MONO, False)]])

# ── 2. Problem ──────────────────────────────────────────────────────────────
s = slide()
kicker(s, 0.9, 0.7, "The problem", PINK)
rule(s, 0.9, 1.2, 3.2, PINK)
text(s, 0.9, 1.5, 11.5, 1.3,
     [[("13 million kirana stores run India. ", 30, WHITE, DISPLAY, True),
       ("They run on paper.", 30, PINK, DISPLAY, True)]])
for i, (big, small) in enumerate([
    ("Invisible", "The day's sales, the udhaar owed, the stock running low live in a notebook under the counter."),
    ("No insight", "She can't see her own profit or cash flow — no P&L, no idea what's bleeding her margin."),
    ("No credit", "A 20-year track record earns her ZERO formal credit. The bank can't see paper."),
    ("Under siege", "Quick-commerce has shut 200,000 kiranas since 2023."),
]):
    y = 3.0 + i * 1.05
    text(s, 0.9, y, 2.6, 0.9, [[(big, 22, MINT, DISPLAY, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 3.7, y, 8.7, 0.9, [[(small, 16, WHITE, BODY, False)]], anchor=MSO_ANCHOR.MIDDLE)

# ── 3. Solution loop ────────────────────────────────────────────────────────
s = slide()
kicker(s, 0.9, 0.7, "The solution · one loop, all on the phone", MINT)
rule(s, 0.9, 1.2, 6.0, UV)
text(s, 0.9, 1.45, 12, 0.8,
     [[("capture → understand → advise → sell → credit", 30, WHITE, DISPLAY, True)]])
steps = [
    ("1 · CAPTURE", MINT, "Photograph a ledger, a customer bill, or a wholesaler challan — or just speak it. Claude Opus 4.8 vision reads Hindi handwriting; whisper.cpp handles voice on-device. Zero typing."),
    ("2 · UNDERSTAND", MINT, "Entries auto-build a live P&L, inventory (cost→sell margins, low-stock), and udhaar tracker. Every line snapshots price+cost — credit-grade data, not estimates."),
    ("3 · ADVISE", YELLOW, "Four AI agents brief her daily; an agentic assistant answers anything in her language — with inline charts — and drafts sales/payments she confirms."),
    ("4 · SELL", MINT, "One tap publishes her live stock as a neighbourhood storefront — her counter-punch to quick commerce, from data she already has."),
    ("5 · CREDIT", UV, "Every rupee feeds an Artha Score (0–100) + working-capital band — an alternative-data credit profile from real commerce flow. The first credit she was ever offered."),
]
for i, (head, col, body) in enumerate(steps):
    y = 2.65 + i * 0.93
    text(s, 0.9, y, 2.5, 0.85, [[(head, 15, col, MONO, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 3.5, y, 9.0, 0.85, [[(body, 14.5, WHITE, BODY, False)]], anchor=MSO_ANCHOR.MIDDLE)

# ── 4 & 5. Product screenshots ──────────────────────────────────────────────
def shots_slide(title, items):
    s = slide()
    kicker(s, 0.9, 0.55, title, MINT)
    rule(s, 0.9, 1.02, 3.0, UV)
    h = 5.1
    w = h * PHONE_RATIO
    gap = 0.7
    total = len(items) * w + (len(items) - 1) * gap
    x = (13.333 - total) / 2
    for name, cap, sub in items:
        phone(s, name, x, 1.35, h)
        text(s, x - 0.25, 6.6, w + 0.5, 0.7,
             [[(cap, 14, MINT, MONO, True)],
              [(sub, 11.5, MUTED, BODY, False)]], align=PP_ALIGN.CENTER)
        x += w + gap
    return s

shots_slide("The product (1/2)", [
    ("home.png",         "HOME + DATE NAV", "Any day's P&L · ◄ date ►"),
    ("pnl.png",          "P&L DASHBOARD",   "Revenue · profit · udhaar · chart"),
    ("scan-sales.png",   "SCAN A SALE",     "Customer bill / day scribble"),
])
shots_slide("The product (2/2)", [
    ("scan-challan.png",     "SCAN A CHALLAN",   "Supplier bill → priced stock"),
    ("assistant.png",        "ASK IN HINDI",     "Voice or text, your language"),
    ("assistant-chart.png",  "AGENT + CHARTS",   "Reads live data, answers visually"),
])

# ── 6. Architecture ─────────────────────────────────────────────────────────
s = slide()
kicker(s, 0.9, 0.65, "Architecture · local-first, on-device edge", MINT)
rule(s, 0.9, 1.13, 5.0, UV)
text(s, 0.9, 1.38, 11.7, 0.9,
     [[("Runs on the iQOO's edge — nothing leaves the phone but an image for OCR.", 22, WHITE, DISPLAY, True)]])
arch = [
    ("ON-DEVICE LLM", "Qwen 2.5 3B via llama.cpp parses entries, routes intents, and runs the insight loop on the iQOO — offline."),
    ("CLOUD = OCR ONLY", "Claude Opus 4.8 vision reads Hindi handwriting — the one task that genuinely needs a frontier model. The sole cloud touchpoint."),
    ("VOICE", "On-device Hindi STT — fine-tuned whisper-hindi-small (whisper.cpp, JNI, arm64)."),
    ("DATA", "Local Room/SQLite v3 — customers, price snapshots, reactive flows. The credit-grade record stays on the phone."),
    ("AGENT", "Tool-calling over 8 read-only shop-data tools + draft-action tools → confirm cards; answers with inline charts."),
    ("APP", "Kotlin · Jetpack Compose · Clean Arch · Hilt · 'The Verge' design system. OfficeKit bridges the build."),
]
for i, (head, body) in enumerate(arch):
    y = 2.45 + i * 0.8
    text(s, 0.9, y, 2.7, 0.75, [[(head, 13.5, MINT, MONO, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 3.7, y, 8.9, 0.75, [[(body, 13.5, WHITE, BODY, False)]], anchor=MSO_ANCHOR.MIDDLE)

# ── 7. Why it holds / close ─────────────────────────────────────────────────
s = slide()
rule(s, 0.9, 1.05, 4.2, UV)
kicker(s, 0.9, 0.6, "Why it holds where others broke")
text(s, 0.9, 1.35, 11.5, 2.2, [
    [("The ledger is ", 22, WHITE, BODY, False), ("free", 22, MINT, BODY, True),
     (" — acquisition + data, never the revenue.", 22, WHITE, BODY, False)],
    [("The storefront builds ", 22, WHITE, BODY, False), ("daily habit", 22, MINT, BODY, True),
     (" and a digital order trail.", 22, WHITE, BODY, False)],
    [("The credit layer is the ", 22, WHITE, BODY, False), ("monetisation", 22, MINT, BODY, True),
     (" the category never had.", 22, WHITE, BODY, False)],
], space=10)
text(s, 0.9, 4.0, 11.5, 1.4,
     [[("Three PS angles — phone-as-POS, credit for the unbanked, a vernacular assistant — ", 18, MUTED, BODY, False),
       ("fused into one loop.", 18, MINT, BODY, True)]])
text(s, 0.9, 5.5, 11.5, 1.2,
     [[("ARTHA", 40, WHITE, DISPLAY, True)],
      [("github.com/Om-Midya/artha-kirana", 16, MINT, MONO, True)]])

out = os.path.join(HERE, "Artha-Pitch.pptx")
prs.save(out)
print("saved", out, "·", len(prs.slides._sldIdLst), "slides")
